#!/usr/bin/env python3
import tkinter as tk
from tkinter import filedialog
import sys
import os
from pathlib import Path
import warnings

# Подавляем предупреждения pandas
warnings.filterwarnings('ignore')

# Дополнительно подавляем конкретные предупреждения
import warnings
warnings.filterwarnings('ignore', category=UserWarning, module='pandas')
warnings.filterwarnings('ignore', category=FutureWarning)

# Проверяем наличие необходимых библиотек для офисных форматов
def check_dependencies():
    missing = []
    try:
        import PyPDF2
    except ImportError:
        missing.append("PyPDF2")
    
    try:
        import docx
    except ImportError:
        missing.append("python-docx")
    
    try:
        import pptx
    except ImportError:
        missing.append("python-pptx")
    
    try:
        import openpyxl
    except ImportError:
        missing.append("openpyxl")
    
    try:
        import pandas as pd
    except ImportError:
        missing.append("pandas")
    
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        missing.append("beautifulsoup4")
    
    return missing

def extract_text_from_pdf(file_path):
    """Извлекает текст из PDF файла"""
    try:
        import PyPDF2
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text.strip()
    except Exception as e:
        return f"[Ошибка при чтении PDF: {str(e)}]"

def extract_text_from_docx(file_path):
    """Извлекает текст из DOCX файла"""
    try:
        import docx
        doc = docx.Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        return f"[Ошибка при чтении DOCX: {str(e)}]"

def extract_text_from_pptx(file_path):
    """Извлекает текст из PPTX файла"""
    try:
        import pptx
        prs = pptx.Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        return f"[Ошибка при чтении PPTX: {str(e)}]"

def extract_text_from_excel(file_path):
    """Извлекает текст из Excel файлов"""
    try:
        import pandas as pd
        import openpyxl
        
        text = ""
        try:
            # Для .xlsx
            xl = pd.ExcelFile(file_path)
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name)
                
                # Преобразуем в текстовый формат с табуляцией
                for idx, row in df.iterrows():
                    row_text = "\t".join([str(cell) for cell in row.values if pd.notna(cell) and str(cell).strip()])
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
        except Exception:
            # Альтернативный способ для старых форматов
            import xlrd
            workbook = xlrd.open_workbook(file_path)
            for sheet in workbook.sheets():
                for row in range(sheet.nrows):
                    row_data = []
                    for col in range(sheet.ncols):
                        cell_value = sheet.cell_value(row, col)
                        if cell_value and str(cell_value).strip():
                            row_data.append(str(cell_value))
                    if row_data:
                        text += "\t".join(row_data) + "\n"
                text += "\n"
        
        return text.strip()
    except Exception as e:
        return f"[Ошибка при чтении Excel: {str(e)}]"

def extract_text_from_html(file_path):
    """Извлекает текст из HTML файла"""
    try:
        from bs4 import BeautifulSoup
        with open(file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file, 'html.parser')
            # Удаляем скрипты и стили
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text()
            # Убираем лишние пробелы
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
        return text.strip()
    except Exception as e:
        return f"[Ошибка при чтении HTML: {str(e)}]"

def extract_text_from_xml(file_path):
    """Извлекает текст из XML файла"""
    try:
        from bs4 import BeautifulSoup
        with open(file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file, 'xml')
            text = soup.get_text()
            # Убираем лишние пробелы
            lines = (line.strip() for line in text.splitlines())
            text = '\n'.join(line for line in lines if line)
        return text.strip()
    except Exception as e:
        return f"[Ошибка при чтении XML: {str(e)}]"

def extract_text_from_json(file_path):
    """Читает JSON файл с сохранением форматирования"""
    try:
        import json
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
            # Форматируем с отступами
            text = json.dumps(data, ensure_ascii=False, indent=2)
        return text.strip()
    except Exception as e:
        return f"[Ошибка при чтении JSON: {str(e)}]"

def extract_text_from_csv(file_path):
    """Читает CSV файл"""
    try:
        import csv
        text = ""
        with open(file_path, 'r', encoding='utf-8') as file:
            csv_reader = csv.reader(file)
            for row in csv_reader:
                if any(cell.strip() for cell in row):
                    text += "\t".join(row) + "\n"
        return text.strip()
    except Exception as e:
        return f"[Ошибка при чтении CSV: {str(e)}]"

def read_text_file_with_encodings(file_path):
    """Читает текстовый файл, пробуя разные кодировки"""
    encodings = ['utf-8', 'cp1251', 'iso-8859-1', 'mac_roman', 'windows-1252']
    
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as file:
                content = file.read()
                return content
        except UnicodeDecodeError:
            continue
        except Exception:
            continue
    
    # Если ни одна кодировка не подошла, пробуем бинарный режим
    try:
        with open(file_path, 'rb') as file:
            return file.read().decode('utf-8', errors='ignore')
    except Exception:
        return f"[Не удалось прочитать файл {file_path}]"

# Словарь с обработчиками для разных форматов
FILE_HANDLERS = {
    # Программистские форматы
    '.py': read_text_file_with_encodings,
    '.js': read_text_file_with_encodings,
    '.java': read_text_file_with_encodings,
    '.cpp': read_text_file_with_encodings,
    '.c': read_text_file_with_encodings,
    '.cs': read_text_file_with_encodings,
    '.php': read_text_file_with_encodings,
    '.rb': read_text_file_with_encodings,
    '.go': read_text_file_with_encodings,
    '.rs': read_text_file_with_encodings,
    '.swift': read_text_file_with_encodings,
    '.kt': read_text_file_with_encodings,
    '.scala': read_text_file_with_encodings,
    '.pl': read_text_file_with_encodings,
    '.lua': read_text_file_with_encodings,
    '.sql': read_text_file_with_encodings,
    '.r': read_text_file_with_encodings,
    '.m': read_text_file_with_encodings,  # MATLAB/Objective-C
    '.f': read_text_file_with_encodings,  # Fortran
    '.pas': read_text_file_with_encodings,  # Pascal
    '.vb': read_text_file_with_encodings,  # Visual Basic
    
    # Веб-разработка
    '.html': extract_text_from_html,
    '.htm': extract_text_from_html,
    '.css': read_text_file_with_encodings,
    '.scss': read_text_file_with_encodings,
    '.sass': read_text_file_with_encodings,
    '.less': read_text_file_with_encodings,
    '.jsx': read_text_file_with_encodings,
    '.ts': read_text_file_with_encodings,
    '.tsx': read_text_file_with_encodings,
    '.vue': read_text_file_with_encodings,
    
    # Конфигурационные файлы
    '.xml': extract_text_from_xml,
    '.json': extract_text_from_json,
    '.yaml': read_text_file_with_encodings,
    '.yml': read_text_file_with_encodings,
    '.toml': read_text_file_with_encodings,
    '.ini': read_text_file_with_encodings,
    '.cfg': read_text_file_with_encodings,
    '.conf': read_text_file_with_encodings,
    '.properties': read_text_file_with_encodings,
    
    # Разметка и документация
    '.md': read_text_file_with_encodings,
    '.markdown': read_text_file_with_encodings,
    '.rst': read_text_file_with_encodings,
    '.tex': read_text_file_with_encodings,
    
    # Скрипты
    '.sh': read_text_file_with_encodings,
    '.bash': read_text_file_with_encodings,
    '.zsh': read_text_file_with_encodings,
    '.ps1': read_text_file_with_encodings,
    '.bat': read_text_file_with_encodings,
    '.cmd': read_text_file_with_encodings,
    
    # Данные
    '.csv': extract_text_from_csv,
    '.tsv': extract_text_from_csv,
    
    # Офисные форматы
    '.pdf': extract_text_from_pdf,
    '.docx': extract_text_from_docx,
    '.doc': extract_text_from_docx,
    '.pptx': extract_text_from_pptx,
    '.ppt': extract_text_from_pptx,
    '.xlsx': extract_text_from_excel,
    '.xls': extract_text_from_excel,
    '.xlsm': extract_text_from_excel,
    '.odt': read_text_file_with_encodings,
    '.ods': read_text_file_with_encodings,
    '.odp': read_text_file_with_encodings,
    
    # Текстовые файлы
    '.txt': read_text_file_with_encodings,
    '.log': read_text_file_with_encodings,
    '.rtf': read_text_file_with_encodings,
}

def convert_files_to_txt(file_paths, output_dir):
    """Конвертирует выбранные файлы в текстовые версии"""
    
    # Получаем список отсутствующих зависимостей
    missing_deps = check_dependencies()
    if missing_deps:
        print("Внимание: отсутствуют библиотеки для некоторых форматов:")
        for dep in missing_deps:
            print(f"  - {dep}")
        print("\nУстановите командой: pip install " + " ".join(missing_deps))
        print("Файлы в этих форматах могут не конвертироваться правильно.\n")
    
    converted_count = 0
    skipped_count = 0
    
    for file_path in file_paths:
        try:
            file_ext = Path(file_path).suffix.lower()
            
            # Проверяем, поддерживается ли формат
            if file_ext not in FILE_HANDLERS:
                print(f"⚠ Неподдерживаемый формат: {file_ext} ({Path(file_path).name})")
                skipped_count += 1
                continue
            
            # Получаем текст из файла
            handler = FILE_HANDLERS[file_ext]
            text_content = handler(file_path)
            
            if text_content.startswith("[Ошибка"):
                print(f"✗ Ошибка: {Path(file_path).name} - {text_content}")
                skipped_count += 1
                continue
            
            # Создаем имя для нового файла
            original_name = Path(file_path).stem
            new_filename = f"{original_name}.txt"
            
            # Заменяем недопустимые символы в имени файла
            invalid_chars = '<>:"/\\|?*'
            for char in invalid_chars:
                new_filename = new_filename.replace(char, '_')
            
            # Формируем полный путь для сохранения
            output_path = Path(output_dir) / new_filename
            
            # Если файл уже существует, добавляем номер
            counter = 1
            while output_path.exists():
                new_filename = f"{original_name}_{counter}.txt"
                output_path = Path(output_dir) / new_filename
                counter += 1
            
            # Сохраняем текст в файл с кодировкой UTF-8
            with open(output_path, 'w', encoding='utf-8') as dest_file:
                dest_file.write(text_content)
            
            print(f"✓ Конвертирован: {Path(file_path).name} -> {output_path.name}")
            converted_count += 1
            
        except Exception as e:
            print(f"✗ Ошибка конвертации {Path(file_path).name}: {str(e)}")
            skipped_count += 1
    
    return converted_count, skipped_count

def main():
    # Определяем путь для сохранения
    if len(sys.argv) > 1:
        output_dir = sys.argv[1]
    else:
        output_dir = "/Users/aleksej/Downloads"
    
    # Создаем директорию если её нет
    os.makedirs(output_dir, exist_ok=True)
    
    # Настраиваем и запускаем Tkinter для выбора файлов
    root = tk.Tk()
    root.withdraw()  # Скрываем основное окно
    
    # Делаем окно поверх всех окон
    root.attributes('-topmost', 1)
    root.update()
    
    # Открываем диалог выбора файлов
    file_paths = filedialog.askopenfilenames(
        title="Выберите файлы для конвертации",
        filetypes=[
            ("Все файлы", "*.*"),
            ("Программистские файлы", "*.py *.js *.java *.cpp *.c *.cs *.php *.rb *.go"),
            ("Веб-файлы", "*.html *.htm *.css *.jsx *.ts *.tsx"),
            ("Офисные файлы", "*.pdf *.docx *.doc *.pptx *.ppt *.xlsx *.xls"),
            ("Конфигурационные файлы", "*.json *.xml *.yaml *.yml *.ini"),
            ("Текстовые файлы", "*.txt *.md *.log"),
        ]
    )
    
    # Возвращаем обычное поведение окна
    root.attributes('-topmost', 0)
    
    if not file_paths:
        print("Файлы не выбраны")
        return
    
    print(f"Выбрано файлов: {len(file_paths)}")
    print(f"Путь сохранения: {output_dir}")
    
    # Конвертируем выбранные файлы
    converted, skipped = convert_files_to_txt(file_paths, output_dir)
    
    print(f"Готово!")
    print(f"Успешно конвертировано: {converted}")
    print(f"Пропущено: {skipped}")
    print(f"Файлы сохранены в: {output_dir}")
    
    # Предлагаем открыть папку с результатами
    if converted > 0:
        answer = input("\nОткрыть папку с результатами? (y/n): ")
        if answer.lower() in ['y', 'yes', 'д', 'да']:
            os.system(f'open "{output_dir}"')

if __name__ == "__main__":
    main()